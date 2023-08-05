# -*- coding: utf-8 -*-
#
# test_shapes.py
#
# Copyright (C) 2012, 2013 Steve Canny scanny@cisco.com
#
# This module is part of python-pptx and is released under
# the MIT License: http://www.opensource.org/licenses/mit-license.php

"""Test suite for pptx.shapes module."""

import os

from hamcrest import assert_that, equal_to, is_

from mock import Mock, patch, PropertyMock

from pptx.constants import MSO
from pptx.oxml import (
    _SubElement, nsdecls, oxml_fromstring, oxml_parse, oxml_tostring
)
from pptx.presentation import _SlideLayout
from pptx.shapes import (
    _AutoShapeType, _BaseShape, _Cell, _CellCollection, _Column,
    _ColumnCollection, _Font, _Paragraph, _Placeholder, _Row, _RowCollection,
    _Run, _Shape, _ShapeCollection, _TextFrame, _to_unicode
)
from pptx.spec import namespaces
from pptx.spec import (
    PH_TYPE_CTRTITLE, PH_TYPE_DT, PH_TYPE_FTR, PH_TYPE_OBJ, PH_TYPE_SLDNUM,
    PH_TYPE_SUBTITLE, PH_TYPE_TBL, PH_TYPE_TITLE, PH_ORIENT_HORZ,
    PH_ORIENT_VERT, PH_SZ_FULL, PH_SZ_HALF, PH_SZ_QUARTER
)
from pptx.util import Inches, Pt
from testdata import test_shape_elements, test_shapes
from testing import TestCase


# module globals -------------------------------------------------------------
def absjoin(*paths):
    return os.path.abspath(os.path.join(*paths))

thisdir = os.path.split(__file__)[0]
test_file_dir = absjoin(thisdir, 'test_files')

test_image_path = absjoin(test_file_dir, 'python-icon.jpeg')
test_bmp_path = absjoin(test_file_dir, 'python.bmp')
new_image_path = absjoin(test_file_dir, 'monty-truth.png')
test_pptx_path = absjoin(test_file_dir, 'test.pptx')
images_pptx_path = absjoin(test_file_dir, 'with_images.pptx')

nsmap = namespaces('a', 'r', 'p')


def _sldLayout1():
    path = os.path.join(thisdir, 'test_files/slideLayout1.xml')
    sldLayout = oxml_parse(path).getroot()
    return sldLayout


def _sldLayout1_shapes():
    sldLayout = _sldLayout1()
    spTree = sldLayout.xpath('./p:cSld/p:spTree', namespaces=nsmap)[0]
    shapes = _ShapeCollection(spTree)
    return shapes


class Test_AutoShapeType(TestCase):
    """Test _AutoShapeType"""
    def test_construction_return_values(self):
        """_AutoShapeType() returns instance with correct property values"""
        # setup ------------------------
        id_ = MSO.SHAPE_ROUNDED_RECTANGLE
        prst = 'roundRect'
        basename = 'Rounded Rectangle'
        desc = 'Rounded rectangle'
        # exercise ---------------------
        autoshape_type = _AutoShapeType(id_)
        # verify -----------------------
        assert_that(autoshape_type.autoshape_type_id, is_(equal_to(id_)))
        assert_that(autoshape_type.prst, is_(equal_to(prst)))
        assert_that(autoshape_type.basename, is_(equal_to(basename)))
        assert_that(autoshape_type.desc, is_(equal_to(desc)))

    def test__lookup_id_by_prst_return_value(self):
        """_AutoShapeType._lookup_id_by_prst() return value is correct"""
        # setup ------------------------
        autoshape_type_id = MSO.SHAPE_ROUNDED_RECTANGLE
        prst = 'roundRect'
        # exercise ---------------------
        retval = _AutoShapeType._lookup_id_by_prst(prst)
        # verify -----------------------
        assert_that(retval, is_(equal_to(autoshape_type_id)))

    def test__lookup_id_raises_on_bad_prst(self):
        """_AutoShapeType._lookup_id_by_prst() raises on bad prst"""
        # setup ------------------------
        prst = 'badPrst'
        # verify -----------------------
        with self.assertRaises(KeyError):
            _AutoShapeType._lookup_id_by_prst(prst)

    def test_second_construction_returns_cached_instance(self):
        """_AutoShapeType() returns cached instance on duplicate call"""
        # setup ------------------------
        id_ = MSO.SHAPE_ROUNDED_RECTANGLE
        ast1 = _AutoShapeType(id_)
        # exercise ---------------------
        ast2 = _AutoShapeType(id_)
        # verify -----------------------
        assert_that(ast2, is_(equal_to(ast1)))

    def test_construction_raises_on_bad_autoshape_type_id(self):
        """_AutoShapeType() raises on bad auto shape type id"""
        # setup ------------------------
        autoshape_type_id = 9999
        # verify -----------------------
        with self.assertRaises(KeyError):
            _AutoShapeType(autoshape_type_id)


class Test_BaseShape(TestCase):
    """Test _BaseShape"""
    def setUp(self):
        path = os.path.join(thisdir, 'test_files/slide1.xml')
        self.sld = oxml_parse(path).getroot()
        xpath = './p:cSld/p:spTree/p:pic'
        pic = self.sld.xpath(xpath, namespaces=nsmap)[0]
        self.base_shape = _BaseShape(pic)

    def test_has_textframe_value(self):
        """_BaseShape.has_textframe value correct"""
        # setup ------------------------
        spTree = self.sld.xpath('./p:cSld/p:spTree', namespaces=nsmap)[0]
        shapes = _ShapeCollection(spTree)
        indexes = []
        # exercise ---------------------
        for idx, shape in enumerate(shapes):
            if shape.has_textframe:
                indexes.append(idx)
        # verify -----------------------
        expected = [0, 1, 3, 5, 6]
        actual = indexes
        msg = ("expected txBody element in shapes %s, got %s" %
               (expected, actual))
        self.assertEqual(expected, actual, msg)

    def test_id_value(self):
        """_BaseShape.id value is correct"""
        # exercise ---------------------
        id = self.base_shape.id
        # verify -----------------------
        expected = 6
        actual = id
        msg = "expected %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_is_placeholder_true_for_placeholder(self):
        """_BaseShape.is_placeholder True for placeholder shape"""
        # setup ------------------------
        xpath = './p:cSld/p:spTree/p:sp'
        sp = self.sld.xpath(xpath, namespaces=nsmap)[0]
        base_shape = _BaseShape(sp)
        # verify -----------------------
        actual = base_shape.is_placeholder
        msg = "expected True, got %s" % (actual)
        self.assertTrue(actual, msg)

    def test_is_placeholder_false_for_non_placeholder(self):
        """_BaseShape.is_placeholder False for non-placeholder shape"""
        # verify -----------------------
        actual = self.base_shape.is_placeholder
        msg = "expected False, got %s" % (actual)
        self.assertFalse(actual, msg)

    def test__is_title_true_for_title_placeholder(self):
        """_BaseShape._is_title True for title placeholder shape"""
        # setup ------------------------
        xpath = './p:cSld/p:spTree/p:sp'
        title_placeholder_sp = self.sld.xpath(xpath, namespaces=nsmap)[0]
        base_shape = _BaseShape(title_placeholder_sp)
        # verify -----------------------
        actual = base_shape._is_title
        msg = "expected True, got %s" % (actual)
        self.assertTrue(actual, msg)

    def test__is_title_false_for_no_ph_element(self):
        """_BaseShape._is_title False on shape has no <p:ph> element"""
        # setup ------------------------
        self.base_shape._element = Mock(name='_element')
        self.base_shape._element.xpath.return_value = []
        # verify -----------------------
        assert_that(self.base_shape._is_title, is_(False))

    def test_name_value(self):
        """_BaseShape.name value is correct"""
        # exercise ---------------------
        name = self.base_shape.name
        # verify -----------------------
        expected = 'Picture 5'
        actual = name
        msg = "expected '%s', got '%s'" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_shape_name_returns_none_for_unimplemented_shape_types(self):
        """_BaseShape.shape_name returns None for unimplemented shape types"""
        assert_that(self.base_shape.shape_type, is_(None))

    def test_textframe_raises_on_no_textframe(self):
        """_BaseShape.textframe raises on shape with no text frame"""
        with self.assertRaises(ValueError):
            self.base_shape.textframe

    def test_text_setter_structure_and_value(self):
        """assign to _BaseShape.text yields single run para set to value"""
        # setup ------------------------
        test_text = 'python-pptx was here!!'
        xpath = './p:cSld/p:spTree/p:sp'
        textbox_sp = self.sld.xpath(xpath, namespaces=nsmap)[2]
        base_shape = _BaseShape(textbox_sp)
        # exercise ---------------------
        base_shape.text = test_text
        # verify paragraph count ------
        expected = 1
        actual = len(base_shape.textframe.paragraphs)
        msg = "expected paragraph count %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)
        # verify value ----------------
        expected = test_text
        actual = base_shape.textframe.paragraphs[0].runs[0].text
        msg = "expected text '%s', got '%s'" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_text_setter_raises_on_no_textframe(self):
        """assignment to _BaseShape.text raises for shape with no text frame"""
        with self.assertRaises(TypeError):
            self.base_shape.text = 'test text'


class Test_Cell(TestCase):
    """Test _Cell"""
    def setUp(self):
        tc_xml = '<a:tc %s><a:txBody><a:p/></a:txBody></a:tc>' % nsdecls('a')
        test_tc_elm = oxml_fromstring(tc_xml)
        self.cell = _Cell(test_tc_elm)

    def test_text_round_trips_intact(self):
        """_Cell.text (setter) sets cell text"""
        # setup ------------------------
        test_text = 'test_text'
        # exercise ---------------------
        self.cell.text = test_text
        # verify -----------------------
        text = self.cell.textframe.paragraphs[0].runs[0].text
        assert_that(text, is_(equal_to(test_text)))


class Test_CellCollection(TestCase):
    """Test _CellCollection"""
    def setUp(self):
        tr_xml = (
            '<a:tr %s h="370840"><a:tc><a:txBody><a:p/></a:txBody></a:tc><a:t'
            'c><a:txBody><a:p/></a:txBody></a:tc></a:tr>' % nsdecls('a')
        )
        test_tr_elm = oxml_fromstring(tr_xml)
        self.cells = _CellCollection(test_tr_elm)

    def test_is_indexable(self):
        """_CellCollection indexable (e.g. no TypeError on 'cells[0]')"""
        # verify -----------------------
        try:
            self.cells[0]
        except TypeError:
            msg = "'_CellCollection' object does not support indexing"
            self.fail(msg)
        except IndexError:
            pass

    def test_is_iterable(self):
        """_CellCollection is iterable (e.g. ``for cell in cells:``)"""
        # verify -----------------------
        count = 0
        try:
            for cell in self.cells:
                count += 1
        except TypeError:
            msg = "_CellCollection object is not iterable"
            self.fail(msg)
        assert_that(count, is_(equal_to(2)))

    def test_raises_on_idx_out_of_range(self):
        """_CellCollection raises on index out of range"""
        with self.assertRaises(IndexError):
            self.cells[9]

    def test_cell_count_correct(self):
        """len(_CellCollection) returns correct cell count"""
        # verify -----------------------
        assert_that(len(self.cells), is_(equal_to(2)))


class Test_Column(TestCase):
    """Test _Column"""
    def setUp(self):
        gridCol_xml = '<a:gridCol %s w="3048000"/>' % nsdecls('a')
        test_gridCol_elm = oxml_fromstring(gridCol_xml)
        self.column = _Column(test_gridCol_elm, Mock(name='table'))

    def test_width_from_xml_correct(self):
        """_Column.width returns correct value from gridCol XML element"""
        # verify -----------------------
        assert_that(self.column.width, is_(equal_to(3048000)))

    def test_width_round_trips_intact(self):
        """_Column.width round-trips intact"""
        # setup ------------------------
        self.column.width = 999
        # verify -----------------------
        assert_that(self.column.width, is_(equal_to(999)))

    def test_set_width_raises_on_bad_value(self):
        """_Column.width raises on attempt to assign invalid value"""
        test_cases = ('abc', '1', -1)
        for value in test_cases:
            with self.assertRaises(ValueError):
                self.column.width = value


class Test_ColumnCollection(TestCase):
    """Test _ColumnCollection"""
    def setUp(self):
        tbl_xml = (
            '<a:tbl %s><a:tblGrid><a:gridCol w="3048000"/><a:gridCol w="30480'
            '00"/></a:tblGrid></a:tbl>' % nsdecls('a')
        )
        test_tbl_elm = oxml_fromstring(tbl_xml)
        self.columns = _ColumnCollection(test_tbl_elm, Mock(name='table'))

    def test_is_indexable(self):
        """_ColumnCollection indexable (e.g. no TypeError on 'columns[0]')"""
        # verify -----------------------
        try:
            self.columns[0]
        except TypeError:
            msg = "'_ColumnCollection' object does not support indexing"
            self.fail(msg)
        except IndexError:
            pass

    def test_is_iterable(self):
        """_ColumnCollection is iterable (e.g. ``for col in columns:``)"""
        # verify -----------------------
        count = 0
        try:
            for col in self.columns:
                count += 1
        except TypeError:
            msg = "_ColumnCollection object is not iterable"
            self.fail(msg)
        assert_that(count, is_(equal_to(2)))

    def test_raises_on_idx_out_of_range(self):
        """_ColumnCollection raises on index out of range"""
        with self.assertRaises(IndexError):
            self.columns[9]

    def test_column_count_correct(self):
        """len(_ColumnCollection) returns correct column count"""
        # verify -----------------------
        assert_that(len(self.columns), is_(equal_to(2)))


class Test_Font(TestCase):
    """Test _Font class"""
    def setUp(self):
        self.rPr_xml = '<a:rPr %s/>' % nsdecls('a')
        self.rPr = oxml_fromstring(self.rPr_xml)
        self.font = _Font(self.rPr)

    def test_get_bold_setting(self):
        """_Font.bold returns True on bold font weight"""
        # setup ------------------------
        rPr_xml = '<a:rPr %s b="1"/>' % nsdecls('a')
        rPr = oxml_fromstring(rPr_xml)
        font = _Font(rPr)
        # verify -----------------------
        assert_that(self.font.bold, is_(False))
        assert_that(font.bold, is_(True))

    def test_set_bold(self):
        """Setting _Font.bold to True selects bold font weight"""
        # setup ------------------------
        expected_rPr_xml = (
            '<a:rPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006'
            '/main" b="1"/>')
        # exercise ---------------------
        self.font.bold = True
        # verify -----------------------
        rPr_xml = oxml_tostring(self.font._Font__rPr)
        assert_that(rPr_xml, is_(equal_to(expected_rPr_xml)))

    def test_clear_bold(self):
        """Setting _Font.bold to False selects normal font weight"""
        # setup ------------------------
        rPr_xml = (
            '<a:rPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006'
            '/main" b="1"/>')
        rPr = oxml_fromstring(rPr_xml)
        font = _Font(rPr)
        expected_rPr_xml = (
            '<a:rPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006'
            '/main"/>')
        # exercise ---------------------
        font.bold = False
        # verify -----------------------
        rPr_xml = oxml_tostring(font._Font__rPr)
        assert_that(rPr_xml, is_(equal_to(expected_rPr_xml)))

    def test_set_font_size(self):
        """Assignment to _Font.size changes font size"""
        # setup ------------------------
        newfontsize = 2400
        expected_xml = (
            '<a:rPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006'
            '/main" sz="%d"/>') % newfontsize
        # exercise ---------------------
        self.font.size = newfontsize
        # verify -----------------------
        actual_xml = oxml_tostring(self.font._Font__rPr)
        assert_that(actual_xml, is_(equal_to(expected_xml)))


class Test_Paragraph(TestCase):
    """Test _Paragraph"""
    def setUp(self):
        path = os.path.join(thisdir, 'test_files/slide1.xml')
        self.sld = oxml_parse(path).getroot()
        xpath = './p:cSld/p:spTree/p:sp/p:txBody/a:p'
        self.pList = self.sld.xpath(xpath, namespaces=nsmap)

        self.test_text = 'test text'
        self.p_xml = ('<a:p %s><a:r><a:t>%s</a:t></a:r></a:p>' %
                      (nsdecls('a'), self.test_text))
        self.p = oxml_fromstring(self.p_xml)
        self.paragraph = _Paragraph(self.p)

    def test_runs_size(self):
        """_Paragraph.runs is expected size"""
        # setup ------------------------
        actual_lengths = []
        for p in self.pList:
            paragraph = _Paragraph(p)
            # exercise ----------------
            actual_lengths.append(len(paragraph.runs))
        # verify ------------------
        expected = [0, 0, 2, 1, 1, 1]
        actual = actual_lengths
        msg = "expected run count %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_add_run_increments_run_count(self):
        """_Paragraph.add_run() increments run count"""
        # setup ------------------------
        p_elm = self.pList[0]
        paragraph = _Paragraph(p_elm)
        # exercise ---------------------
        paragraph.add_run()
        # verify -----------------------
        expected = 1
        actual = len(paragraph.runs)
        msg = "expected run count %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_clear_removes_all_runs(self):
        """_Paragraph.clear() removes all runs from paragraph"""
        # setup ------------------------
        p = self.pList[2]
        _SubElement(p, 'a:pPr')
        paragraph = _Paragraph(p)
        assert_that(len(paragraph.runs), is_(equal_to(2)))
        # exercise ---------------------
        paragraph.clear()
        # verify -----------------------
        assert_that(len(paragraph.runs), is_(equal_to(0)))

    def test_clear_preserves_paragraph_properties(self):
        """_Paragraph.clear() preserves paragraph properties"""
        # setup ------------------------
        p_xml = ('<a:p %s><a:pPr lvl="1"/><a:r><a:t>%s</a:t></a:r></a:p>' %
                 (nsdecls('a'), self.test_text))
        p_elm = oxml_fromstring(p_xml)
        paragraph = _Paragraph(p_elm)
        expected_p_xml = '<a:p %s><a:pPr lvl="1"/></a:p>' % nsdecls('a')
        # exercise ---------------------
        paragraph.clear()
        # verify -----------------------
        p_xml = oxml_tostring(paragraph._Paragraph__p)
        assert_that(p_xml, is_(equal_to(expected_p_xml)))

    def test_level_setter_generates_correct_xml(self):
        """_Paragraph.level setter generates correct XML"""
        # setup ------------------------
        expected_xml = (
            '<a:p %s>\n  <a:pPr lvl="2"/>\n  <a:r>\n    <a:t>test text</a:t>'
            '\n  </a:r>\n</a:p>\n' % nsdecls('a')
        )
        # exercise ---------------------
        self.paragraph.level = 2
        # verify -----------------------
        self.assertEqualLineByLine(expected_xml, self.paragraph._Paragraph__p)

    def test_level_default_is_zero(self):
        """_Paragraph.level defaults to zero on no lvl attribute"""
        # verify -----------------------
        assert_that(self.paragraph.level, is_(equal_to(0)))

    def test_level_roundtrips_intact(self):
        """_Paragraph.level property round-trips intact"""
        # exercise ---------------------
        self.paragraph.level = 5
        # verify -----------------------
        assert_that(self.paragraph.level, is_(equal_to(5)))

    def test_level_raises_on_bad_value(self):
        """_Paragraph.level raises on attempt to assign invalid value"""
        test_cases = ('0', -1, 9)
        for value in test_cases:
            with self.assertRaises(ValueError):
                self.paragraph.level = value

    def test_set_font_size(self):
        """Assignment to _Paragraph.font.size changes font size"""
        # setup ------------------------
        newfontsize = Pt(54.3)
        expected_xml = (
            '<a:p %s>\n  <a:pPr>\n    <a:defRPr sz="5430"/>\n  </a:pPr>\n  <a'
            ':r>\n    <a:t>test text</a:t>\n  </a:r>\n</a:p>\n' % nsdecls('a')
        )
        # exercise ---------------------
        self.paragraph.font.size = newfontsize
        # verify -----------------------
        self.assertEqualLineByLine(expected_xml, self.paragraph._Paragraph__p)

    def test_text_setter_sets_single_run_text(self):
        """assignment to _Paragraph.text creates single run containing value"""
        # setup ------------------------
        test_text = 'python-pptx was here!!'
        p_elm = self.pList[2]
        paragraph = _Paragraph(p_elm)
        # exercise ---------------------
        paragraph.text = test_text
        # verify -----------------------
        assert_that(len(paragraph.runs), is_(equal_to(1)))
        assert_that(paragraph.runs[0].text, is_(equal_to(test_text)))

    def test_text_accepts_non_ascii_strings(self):
        """assignment of non-ASCII string to text does not raise"""
        # setup ------------------------
        _7bit_string = 'String containing only 7-bit (ASCII) characters'
        _8bit_string = '8-bit string: Hér er texti með íslenskum stöfum.'
        _utf8_literal = u'unicode literal: Hér er texti með íslenskum stöfum.'
        _utf8_from_8bit = unicode('utf-8 unicode: Hér er texti', 'utf-8')
        # verify -----------------------
        try:
            text = _7bit_string
            self.paragraph.text = text
            text = _8bit_string
            self.paragraph.text = text
            text = _utf8_literal
            self.paragraph.text = text
            text = _utf8_from_8bit
            self.paragraph.text = text
        except ValueError:
            msg = "_Paragraph.text rejects valid text string '%s'" % text
            self.fail(msg)


class Test_Placeholder(TestCase):
    """Test _Placeholder"""
    def test_property_values(self):
        """_Placeholder property values are correct"""
        # setup ------------------------
        expected_values = (
            (PH_TYPE_CTRTITLE, PH_ORIENT_HORZ, PH_SZ_FULL,     0),
            (PH_TYPE_DT,       PH_ORIENT_HORZ, PH_SZ_HALF,    10),
            (PH_TYPE_SUBTITLE, PH_ORIENT_VERT, PH_SZ_FULL,     1),
            (PH_TYPE_TBL,      PH_ORIENT_HORZ, PH_SZ_QUARTER, 14),
            (PH_TYPE_SLDNUM,   PH_ORIENT_HORZ, PH_SZ_QUARTER, 12),
            (PH_TYPE_FTR,      PH_ORIENT_HORZ, PH_SZ_QUARTER, 11))
        shapes = _sldLayout1_shapes()
        # exercise ---------------------
        for idx, sp in enumerate(shapes):
            ph = _Placeholder(sp)
            values = (ph.type, ph.orient, ph.sz, ph.idx)
            # verify ----------------------
            expected = expected_values[idx]
            actual = values
            msg = ("expected shapes[%d] values %s, got %s"
                   % (idx, expected, actual))
            self.assertEqual(expected, actual, msg)


class Test_Picture(TestCase):
    """Test _Picture"""
    def test_shape_type_value_correct_for_picture(self):
        """_Shape.shape_type value is correct for picture"""
        # setup ------------------------
        picture = test_shapes.picture
        # verify -----------------------
        assert_that(picture.shape_type, is_(equal_to(MSO.PICTURE)))


class Test_Row(TestCase):
    """Test _Row"""
    def setUp(self):
        tr_xml = (
            '<a:tr %s h="370840"><a:tc><a:txBody><a:p/></a:txBody></a:tc><a:t'
            'c><a:txBody><a:p/></a:txBody></a:tc></a:tr>' % nsdecls('a')
        )
        test_tr_elm = oxml_fromstring(tr_xml)
        self.row = _Row(test_tr_elm, Mock(name='table'))

    def test_height_from_xml_correct(self):
        """_Row.height returns correct value from tr XML element"""
        # verify -----------------------
        assert_that(self.row.height, is_(equal_to(370840)))

    def test_height_round_trips_intact(self):
        """_Row.height round-trips intact"""
        # setup ------------------------
        self.row.height = 999
        # verify -----------------------
        assert_that(self.row.height, is_(equal_to(999)))

    def test_set_height_raises_on_bad_value(self):
        """_Row.height raises on attempt to assign invalid value"""
        test_cases = ('abc', '1', -1)
        for value in test_cases:
            with self.assertRaises(ValueError):
                self.row.height = value


class Test_RowCollection(TestCase):
    """Test _RowCollection"""
    def setUp(self):
        tbl_xml = (
            '<a:tbl %s><a:tr h="370840"><a:tc><a:txBody><a:p/></a:txBody></a:'
            'tc><a:tc><a:txBody><a:p/></a:txBody></a:tc></a:tr><a:tr h="37084'
            '0"><a:tc><a:txBody><a:p/></a:txBody></a:tc><a:tc><a:txBody><a:p/'
            '></a:txBody></a:tc></a:tr></a:tbl>' % nsdecls('a')
        )
        test_tbl_elm = oxml_fromstring(tbl_xml)
        self.rows = _RowCollection(test_tbl_elm, Mock(name='table'))

    def test_is_indexable(self):
        """_RowCollection indexable (e.g. no TypeError on 'rows[0]')"""
        # verify -----------------------
        try:
            self.rows[0]
        except TypeError:
            msg = "'_RowCollection' object does not support indexing"
            self.fail(msg)
        except IndexError:
            pass

    def test_is_iterable(self):
        """_RowCollection is iterable (e.g. ``for row in rows:``)"""
        # verify -----------------------
        count = 0
        try:
            for row in self.rows:
                count += 1
        except TypeError:
            msg = "_RowCollection object is not iterable"
            self.fail(msg)
        assert_that(count, is_(equal_to(2)))

    def test_raises_on_idx_out_of_range(self):
        """_RowCollection raises on index out of range"""
        with self.assertRaises(IndexError):
            self.rows[9]

    def test_row_count_correct(self):
        """len(_RowCollection) returns correct row count"""
        # verify -----------------------
        assert_that(len(self.rows), is_(equal_to(2)))


class Test_Run(TestCase):
    """Test _Run"""
    def setUp(self):
        self.test_text = 'test text'
        self.r_xml = ('<a:r %s><a:t>%s</a:t></a:r>' %
                      (nsdecls('a'), self.test_text))
        self.r = oxml_fromstring(self.r_xml)
        self.run = _Run(self.r)

    def test_set_font_size(self):
        """Assignment to _Run.font.size changes font size"""
        # setup ------------------------
        newfontsize = 2400
        expected_xml = (
            '<a:r %s>\n  <a:rPr sz="2400"/>\n  <a:t>test text</a:t>\n</a:r>\n'
            % nsdecls('a')
        )
        # exercise ---------------------
        self.run.font.size = newfontsize
        # verify -----------------------
        self.assertEqualLineByLine(expected_xml, self.run._Run__r)

    def test_text_value(self):
        """_Run.text value is correct"""
        # exercise ---------------------
        text = self.run.text
        # verify -----------------------
        expected = self.test_text
        actual = text
        msg = "expected '%s', got '%s'" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_text_setter(self):
        """_Run.text setter stores passed value"""
        # setup ------------------------
        new_value = 'new string'
        # exercise ---------------------
        self.run.text = new_value
        # verify -----------------------
        expected = new_value
        actual = self.run.text
        msg = "expected '%s', got '%s'" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test__to_unicode_raises_on_non_string(self):
        """_to_unicode(text) raises on *text* not a string"""
        # verify -----------------------
        with self.assertRaises(TypeError):
            _to_unicode(999)


class Test_Shape(TestCase):
    """Test _Shape"""
    def test_auto_shape_type_value_correct(self):
        """_Shape.auto_shape_type value is correct"""
        # setup ------------------------
        rounded_rectangle = test_shapes.rounded_rectangle
        # verify -----------------------
        assert_that(rounded_rectangle.auto_shape_type,
                    is_(equal_to(MSO.SHAPE_ROUNDED_RECTANGLE)))

    def test_auto_shape_type_raises_on_non_auto_shape(self):
        """_Shape.auto_shape_type raises on non auto shape"""
        # setup ------------------------
        textbox = test_shapes.textbox
        # verify -----------------------
        with self.assertRaises(ValueError):
            textbox.auto_shape_type

    def test_shape_type_value_correct(self):
        """_Shape.shape_type value is correct for all sub-types"""
        # setup ------------------------
        autoshape = test_shapes.autoshape
        placeholder = test_shapes.placeholder
        textbox = test_shapes.textbox
        # verify -----------------------
        assert_that(autoshape.shape_type, is_(equal_to(MSO.AUTO_SHAPE)))
        assert_that(placeholder.shape_type, is_(equal_to(MSO.PLACEHOLDER)))
        assert_that(textbox.shape_type, is_(equal_to(MSO.TEXT_BOX)))

    def test_shape_type_raises_on_unrecognized_shape_type(self):
        """_Shape.shape_type raises on unrecognized shape type"""
        # setup ------------------------
        xml = (
            '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/'
            '2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/'
            '2006/main"><p:nvSpPr><p:cNvPr id="9" name="Unknown Shape Type 8"'
            '/><p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:spPr/></p:sp>'
        )
        sp = oxml_fromstring(xml)
        shape = _Shape(sp)
        # verify -----------------------
        with self.assertRaises(NotImplementedError):
            shape.shape_type


class Test_ShapeCollection(TestCase):
    """Test _ShapeCollection"""
    def setUp(self):
        path = absjoin(test_file_dir, 'slide1.xml')
        sld = oxml_parse(path).getroot()
        spTree = sld.xpath('./p:cSld/p:spTree', namespaces=nsmap)[0]
        self.shapes = _ShapeCollection(spTree)

    def test_construction_size(self):
        """_ShapeCollection is expected size after construction"""
        # verify -----------------------
        self.assertLength(self.shapes, 9)

    def test_constructor_raises_on_contentPart_shape(self):
        """_ShapeCollection() raises on contentPart shape"""
        # setup ------------------------
        spTree = test_shape_elements.empty_spTree
        _SubElement(spTree, 'p:contentPart')
        # verify -----------------------
        with self.assertRaises(ValueError):
            _ShapeCollection(spTree)

    @patch('pptx.shapes.CT_Shape')
    @patch('pptx.shapes._Shape')
    @patch('pptx.shapes._ShapeCollection._ShapeCollection__next_shape_id',
           new_callable=PropertyMock)
    @patch('pptx.shapes._AutoShapeType')
    def test_add_shape_collaboration(self, _AutoShapeType, __next_shape_id,
                                     _Shape, CT_Shape):
        """_ShapeCollection.add_shape() calls the right collaborators"""
        # constant values -------------
        basename = 'Rounded Rectangle'
        prst = 'roundRect'
        id_, name = 9, '%s 8' % basename
        left, top, width, height = 111, 222, 333, 444
        autoshape_type_id = MSO.SHAPE_ROUNDED_RECTANGLE
        # setup mockery ---------------
        autoshape_type = Mock(name='autoshape_type')
        autoshape_type.basename = basename
        autoshape_type.prst = prst
        _AutoShapeType.return_value = autoshape_type
        __next_shape_id.return_value = id_
        sp = Mock(name='sp')
        CT_Shape.new_autoshape_sp.return_value = sp
        __spTree = Mock(name='__spTree')
        __shapes = Mock(name='__shapes')
        shapes = test_shapes.empty_shape_collection
        shapes._ShapeCollection__spTree = __spTree
        shapes._ShapeCollection__shapes = __shapes
        shape = Mock('shape')
        _Shape.return_value = shape
        # exercise ---------------------
        retval = shapes.add_shape(autoshape_type_id, left, top, width, height)
        # verify -----------------------
        _AutoShapeType.assert_called_once_with(autoshape_type_id)
        CT_Shape.new_autoshape_sp.assert_called_once_with(
            id_, name, prst, left, top, width, height)
        _Shape.assert_called_once_with(sp)
        __spTree.append.assert_called_once_with(sp)
        __shapes.append.assert_called_once_with(shape)
        assert_that(retval, is_(equal_to(shape)))

    @patch('pptx.shapes._Picture')
    @patch('pptx.shapes.CT_Picture')
    @patch('pptx.shapes._ShapeCollection._ShapeCollection__next_shape_id',
           new_callable=PropertyMock)
    def test_add_picture_collaboration(self, next_shape_id, CT_Picture,
                                       _Picture):
        """_ShapeCollection.add_picture() calls the right collaborators"""
        # constant values -------------
        file = test_image_path
        left, top, width, height = 1, 2, 3, 4
        id_, name, desc = 12, 'Picture 11', 'image1.jpeg'
        rId = 'rId1'
        # setup mockery ---------------
        next_shape_id.return_value = id_
        image = Mock(name='image', _desc=desc)
        image._scale.return_value = width, height
        rel = Mock(name='rel', _rId=rId)
        slide = Mock(name='slide')
        slide._add_image.return_value = image, rel
        __spTree = Mock(name='__spTree')
        __shapes = Mock(name='__shapes')
        shapes = _ShapeCollection(test_shape_elements.empty_spTree, slide)
        shapes._ShapeCollection__spTree = __spTree
        shapes._ShapeCollection__shapes = __shapes
        pic = Mock(name='pic')
        CT_Picture.new_pic.return_value = pic
        picture = Mock(name='picture')
        _Picture.return_value = picture
        # # exercise --------------------
        retval = shapes.add_picture(file, left, top, width, height)
        # verify -----------------------
        shapes._ShapeCollection__slide._add_image.assert_called_once_with(file)
        image._scale.assert_called_once_with(width, height)
        CT_Picture.new_pic.assert_called_once_with(
            id_, name, desc, rId, left, top, width, height)
        __spTree.append.assert_called_once_with(pic)
        _Picture.assert_called_once_with(pic)
        __shapes.append.assert_called_once_with(picture)
        assert_that(retval, is_(equal_to(picture)))

    @patch('pptx.shapes._Table')
    @patch('pptx.shapes.CT_GraphicalObjectFrame')
    @patch('pptx.shapes._ShapeCollection._ShapeCollection__next_shape_id',
           new_callable=PropertyMock)
    def test_add_table_collaboration(
            self, __next_shape_id, CT_GraphicalObjectFrame, _Table):
        """_ShapeCollection.add_table() calls the right collaborators"""
        # constant values -------------
        id_, name = 9, 'Table 8'
        rows, cols = 2, 3
        left, top, width, height = 111, 222, 333, 444
        # setup mockery ---------------
        __next_shape_id.return_value = id_
        graphicFrame = Mock(name='graphicFrame')
        CT_GraphicalObjectFrame.new_table.return_value = graphicFrame
        __spTree = Mock(name='__spTree')
        __shapes = Mock(name='__shapes')
        shapes = test_shapes.empty_shape_collection
        shapes._ShapeCollection__spTree = __spTree
        shapes._ShapeCollection__shapes = __shapes
        table = Mock('table')
        _Table.return_value = table
        # exercise ---------------------
        retval = shapes.add_table(rows, cols, left, top, width, height)
        # verify -----------------------
        __next_shape_id.assert_called_once_with()
        CT_GraphicalObjectFrame.new_table.assert_called_once_with(
            id_, name, rows, cols, left, top, width, height)
        __spTree.append.assert_called_once_with(graphicFrame)
        _Table.assert_called_once_with(graphicFrame)
        __shapes.append.assert_called_once_with(table)
        assert_that(retval, is_(equal_to(table)))

    @patch('pptx.shapes.CT_Shape')
    @patch('pptx.shapes._Shape')
    @patch('pptx.shapes._ShapeCollection._ShapeCollection__next_shape_id',
           new_callable=PropertyMock)
    def test_add_textbox_collaboration(self, __next_shape_id, _Shape,
                                       CT_Shape):
        """_ShapeCollection.add_textbox() calls the right collaborators"""
        # constant values -------------
        id_, name = 9, 'TextBox 8'
        left, top, width, height = 111, 222, 333, 444
        # setup mockery ---------------
        sp = Mock(name='sp')
        shape = Mock('shape')
        __spTree = Mock(name='__spTree')
        shapes = test_shapes.empty_shape_collection
        shapes._ShapeCollection__spTree = __spTree
        __next_shape_id.return_value = id_
        CT_Shape.new_textbox_sp.return_value = sp
        _Shape.return_value = shape
        # exercise ---------------------
        retval = shapes.add_textbox(left, top, width, height)
        # verify -----------------------
        CT_Shape.new_textbox_sp.assert_called_once_with(
            id_, name, left, top, width, height)
        _Shape.assert_called_once_with(sp)
        __spTree.append.assert_called_once_with(sp)
        assert_that(shapes._ShapeCollection__shapes[0], is_(equal_to(shape)))
        assert_that(retval, is_(equal_to(shape)))

    def test_title_value(self):
        """_ShapeCollection.title value is ref to correct shape"""
        # exercise ---------------------
        title_shape = self.shapes.title
        # verify -----------------------
        expected = 0
        actual = self.shapes.index(title_shape)
        msg = "expected shapes[%d], got shapes[%d]" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_title_is_none_on_no_title_placeholder(self):
        """_ShapeCollection.title value is None when no title placeholder"""
        # setup ------------------------
        shapes = test_shapes.empty_shape_collection
        # verify -----------------------
        assert_that(shapes.title, is_(None))

    def test_placeholders_values(self):
        """_ShapeCollection.placeholders values are correct and sorted"""
        # setup ------------------------
        expected_values = (
            ('Title 1',                    PH_TYPE_CTRTITLE,  0),
            ('Vertical Subtitle 2',        PH_TYPE_SUBTITLE,  1),
            ('Date Placeholder 7',         PH_TYPE_DT,       10),
            ('Footer Placeholder 4',       PH_TYPE_FTR,      11),
            ('Slide Number Placeholder 5', PH_TYPE_SLDNUM,   12),
            ('Table Placeholder 3',        PH_TYPE_TBL,      14))
        shapes = _sldLayout1_shapes()
        # exercise ---------------------
        placeholders = shapes.placeholders
        # verify -----------------------
        for idx, ph in enumerate(placeholders):
            values = (ph.name, ph.type, ph.idx)
            expected = expected_values[idx]
            actual = values
            msg = ("expected placeholders[%d] values %s, got %s" %
                   (idx, expected, actual))
            self.assertEqual(expected, actual, msg)

    def test__clone_layout_placeholders_shapes(self):
        """_ShapeCollection._clone_layout_placeholders clones shapes"""
        # setup ------------------------
        expected_values = (
            [2, 'Title 1',             PH_TYPE_CTRTITLE,  0],
            [3, 'Vertical Subtitle 2', PH_TYPE_SUBTITLE,  1],
            [4, 'Table Placeholder 3', PH_TYPE_TBL,      14])
        slidelayout = _SlideLayout()
        slidelayout._shapes = _sldLayout1_shapes()
        shapes = test_shapes.empty_shape_collection
        # exercise ---------------------
        shapes._clone_layout_placeholders(slidelayout)
        # verify -----------------------
        for idx, sp in enumerate(shapes):
            # verify is placeholder ---
            is_placeholder = sp.is_placeholder
            msg = ("expected shapes[%d].is_placeholder == True %r"
                   % (idx, sp))
            self.assertTrue(is_placeholder, msg)
            # verify values -----------
            ph = _Placeholder(sp)
            expected = expected_values[idx]
            actual = [ph.id, ph.name, ph.type, ph.idx]
            msg = ("expected placeholder[%d] values %s, got %s"
                   % (idx, expected, actual))
            self.assertEqual(expected, actual, msg)

    def test___clone_layout_placeholder_values(self):
        """_ShapeCollection.__clone_layout_placeholder() values correct"""
        # setup ------------------------
        layout_shapes = _sldLayout1_shapes()
        layout_ph_shapes = [sp for sp in layout_shapes if sp.is_placeholder]
        shapes = test_shapes.empty_shape_collection
        expected_values = (
            [2, 'Title 1',                    PH_TYPE_CTRTITLE,  0],
            [3, 'Date Placeholder 2',         PH_TYPE_DT,       10],
            [4, 'Vertical Subtitle 3',        PH_TYPE_SUBTITLE,  1],
            [5, 'Table Placeholder 4',        PH_TYPE_TBL,      14],
            [6, 'Slide Number Placeholder 5', PH_TYPE_SLDNUM,   12],
            [7, 'Footer Placeholder 6',       PH_TYPE_FTR,      11])
        # exercise ---------------------
        for idx, layout_ph_sp in enumerate(layout_ph_shapes):
            layout_ph = _Placeholder(layout_ph_sp)
            sp = shapes._ShapeCollection__clone_layout_placeholder(layout_ph)
            # verify ------------------
            ph = _Placeholder(sp)
            expected = expected_values[idx]
            actual = [ph.id, ph.name, ph.type, ph.idx]
            msg = "expected placeholder values %s, got %s" % (expected, actual)
            self.assertEqual(expected, actual, msg)

    def test___clone_layout_placeholder_xml(self):
        """_ShapeCollection.__clone_layout_placeholder() emits correct XML"""
        # setup ------------------------
        layout_shapes = _sldLayout1_shapes()
        layout_ph_shapes = [sp for sp in layout_shapes if sp.is_placeholder]
        shapes = test_shapes.empty_shape_collection
        expected_xml_tmpl = (
            '<p:sp %s>\n  <p:nvSpPr>\n    <p:cNvPr id="%s" name="%s"/>\n    <'
            'p:cNvSpPr>\n      <a:spLocks noGrp="1"/>\n    </p:cNvSpPr>\n    '
            '<p:nvPr>\n      <p:ph type="%s"%s/>\n    </p:nvPr>\n  </p:nvSpPr'
            '>\n  <p:spPr/>\n%s</p:sp>\n' %
            (nsdecls('p', 'a'), '%d', '%s', '%s', '%s', '%s')
        )
        txBody_snippet = (
            '  <p:txBody>\n    <a:bodyPr/>\n    <a:lstStyle/>\n    <a:p/>\n  '
            '</p:txBody>\n')
        expected_values = [
            (2, 'Title 1', PH_TYPE_CTRTITLE, '', txBody_snippet),
            (3, 'Date Placeholder 2', PH_TYPE_DT, ' sz="half" idx="10"', ''),
            (4, 'Vertical Subtitle 3', PH_TYPE_SUBTITLE,
                ' orient="vert" idx="1"', txBody_snippet),
            (5, 'Table Placeholder 4', PH_TYPE_TBL,
                ' sz="quarter" idx="14"', ''),
            (6, 'Slide Number Placeholder 5', PH_TYPE_SLDNUM,
                ' sz="quarter" idx="12"', ''),
            (7, 'Footer Placeholder 6', PH_TYPE_FTR,
                ' sz="quarter" idx="11"', '')]
                    # verify ----------------------
        for idx, layout_ph_sp in enumerate(layout_ph_shapes):
            layout_ph = _Placeholder(layout_ph_sp)
            sp = shapes._ShapeCollection__clone_layout_placeholder(layout_ph)
            ph = _Placeholder(sp)
            expected_xml = expected_xml_tmpl % expected_values[idx]
            self.assertEqualLineByLine(expected_xml, ph._element)

    def test___next_ph_name_return_value(self):
        """
        _ShapeCollection.__next_ph_name() returns correct value

        * basename + 'Placeholder' + num, e.g. 'Table Placeholder 8'
        * numpart of name defaults to id-1, but increments until unique
        * prefix 'Vertical' if orient="vert"

        """
        cases = (
            (PH_TYPE_OBJ,   3, PH_ORIENT_HORZ, 'Content Placeholder 2'),
            (PH_TYPE_TBL,   4, PH_ORIENT_HORZ, 'Table Placeholder 4'),
            (PH_TYPE_TBL,   7, PH_ORIENT_VERT, 'Vertical Table Placeholder 6'),
            (PH_TYPE_TITLE, 2, PH_ORIENT_HORZ, 'Title 2'))
        # setup ------------------------
        shapes = _sldLayout1_shapes()
        for ph_type, id, orient, expected_name in cases:
            # exercise --------------------
            name = shapes._ShapeCollection__next_ph_name(ph_type, id, orient)
            # verify ----------------------
            expected = expected_name
            actual = name
            msg = ("expected placeholder name '%s', got '%s'" %
                   (expected, actual))
            self.assertEqual(expected, actual, msg)

    def test___next_shape_id_value(self):
        """_ShapeCollection.__next_shape_id value is correct"""
        # setup ------------------------
        shapes = _sldLayout1_shapes()
        # exercise ---------------------
        next_id = shapes._ShapeCollection__next_shape_id
        # verify -----------------------
        expected = 4
        actual = next_id
        msg = "expected %d, got %d" % (expected, actual)
        self.assertEqual(expected, actual, msg)


class Test_Table(TestCase):
    """Test _Table"""
    def test_initial_height_divided_evenly_between_rows(self):
        """Table creation height divided evenly between rows"""
        # constant values -------------
        rows = cols = 3
        left = top = Inches(1.0)
        width = Inches(2.0)
        height = 1000
        shapes = test_shapes.empty_shape_collection
        # exercise ---------------------
        table = shapes.add_table(rows, cols, left, top, width, height)
        # verify -----------------------
        assert_that(table.rows[0].height, is_(equal_to(333)))
        assert_that(table.rows[1].height, is_(equal_to(333)))
        assert_that(table.rows[2].height, is_(equal_to(334)))

    def test_initial_width_divided_evenly_between_columns(self):
        """Table creation width divided evenly between columns"""
        # constant values -------------
        rows = cols = 3
        left = top = Inches(1.0)
        width = 1000
        height = Inches(2.0)
        shapes = test_shapes.empty_shape_collection
        # exercise ---------------------
        table = shapes.add_table(rows, cols, left, top, width, height)
        # verify -----------------------
        assert_that(table.columns[0].width, is_(equal_to(333)))
        assert_that(table.columns[1].width, is_(equal_to(333)))
        assert_that(table.columns[2].width, is_(equal_to(334)))

    def test_height_sum_of_row_heights(self):
        """_Table.height is sum of row heights"""
        # constant values -------------
        rows = cols = 2
        left = top = width = height = Inches(2.0)
        # setup ------------------------
        shapes = test_shapes.empty_shape_collection
        tbl = shapes.add_table(rows, cols, left, top, width, height)
        tbl.rows[0].height = 100
        tbl.rows[1].height = 200
        # verify -----------------------
        sum_of_row_heights = 300
        assert_that(tbl.height, is_(equal_to(sum_of_row_heights)))

    def test_width_sum_of_col_widths(self):
        """_Table.width is sum of column widths"""
        # constant values -------------
        rows = cols = 2
        left = top = width = height = Inches(2.0)
        # setup ------------------------
        shapes = test_shapes.empty_shape_collection
        tbl = shapes.add_table(rows, cols, left, top, width, height)
        tbl.columns[0].width = 100
        tbl.columns[1].width = 200
        # verify -----------------------
        sum_of_col_widths = tbl.columns[0].width + tbl.columns[1].width
        assert_that(tbl.width, is_(equal_to(sum_of_col_widths)))


class Test_TextFrame(TestCase):
    """Test _TextFrame"""
    def setUp(self):
        path = os.path.join(thisdir, 'test_files/slide1.xml')
        self.sld = oxml_parse(path).getroot()
        xpath = './p:cSld/p:spTree/p:sp/p:txBody'
        self.txBodyList = self.sld.xpath(xpath, namespaces=nsmap)

    def test_paragraphs_size(self):
        """_TextFrame.paragraphs is expected size"""
        # setup ------------------------
        actual_lengths = []
        for txBody in self.txBodyList:
            textframe = _TextFrame(txBody)
            # exercise ----------------
            actual_lengths.append(len(textframe.paragraphs))
        # verify -----------------------
        expected = [1, 1, 2, 1, 1]
        actual = actual_lengths
        msg = "expected paragraph count %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_add_paragraph_xml(self):
        """_TextFrame.add_paragraph does what it says"""
        # setup ------------------------
        txBody_xml = (
            '<p:txBody %s><a:bodyPr/><a:p><a:r><a:t>Test text</a:t></a:r></a:'
            'p></p:txBody>' % nsdecls('p', 'a')
        )
        expected_xml = (
            '<p:txBody %s><a:bodyPr/><a:p><a:r><a:t>Test text</a:t></a:r></a:'
            'p><a:p/></p:txBody>' % nsdecls('p', 'a')
        )
        txBody = oxml_fromstring(txBody_xml)
        textframe = _TextFrame(txBody)
        # exercise ---------------------
        textframe.add_paragraph()
        # verify -----------------------
        assert_that(len(textframe.paragraphs), is_(equal_to(2)))
        textframe_xml = oxml_tostring(textframe._TextFrame__txBody)
        expected = expected_xml
        actual = textframe_xml
        msg = "\nExpected: '%s'\n\n     Got: '%s'" % (expected, actual)
        if not expected == actual:
            raise AssertionError(msg)

    def test_text_setter_structure_and_value(self):
        """Assignment to _TextFrame.text yields single run para set to value"""
        # setup ------------------------
        test_text = 'python-pptx was here!!'
        txBody = self.txBodyList[2]
        textframe = _TextFrame(txBody)
        # exercise ---------------------
        textframe.text = test_text
        # verify paragraph count -------
        expected = 1
        actual = len(textframe.paragraphs)
        msg = "expected paragraph count %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)
        # verify value -----------------
        expected = test_text
        actual = textframe.paragraphs[0].runs[0].text
        msg = "expected text '%s', got '%s'" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_vertical_anchor_works(self):
        """Assignment to _TextFrame.vertical_anchor sets vert anchor"""
        # setup ------------------------
        txBody_xml = (
            '<p:txBody %s><a:bodyPr/><a:p><a:r><a:t>Test text</a:t></a:r></a:'
            'p></p:txBody>' % nsdecls('p', 'a')
        )
        expected_xml = (
            '<p:txBody %s>\n  <a:bodyPr anchor="ctr"/>\n  <a:p>\n    <a:r>\n '
            '     <a:t>Test text</a:t>\n    </a:r>\n  </a:p>\n</p:txBody>\n' %
            nsdecls('p', 'a')
        )
        txBody = oxml_fromstring(txBody_xml)
        textframe = _TextFrame(txBody)
        # exercise ---------------------
        textframe.vertical_anchor = MSO.ANCHOR_MIDDLE
        # verify -----------------------
        self.assertEqualLineByLine(expected_xml, textframe._TextFrame__txBody)
